"""Text extraction from uploaded documents.

Each extractor returns a list of `ExtractedPage`s so downstream chunking can preserve
page provenance for citations. Image OCR is stubbed here and enabled in Phase 10.
"""

from __future__ import annotations

import os
from dataclasses import dataclass, field


@dataclass
class ExtractedPage:
    text: str
    page: int | None = None


@dataclass
class ExtractedDoc:
    pages: list[ExtractedPage] = field(default_factory=list)
    title: str | None = None
    page_count: int | None = None

    @property
    def full_text(self) -> str:
        return "\n\n".join(p.text for p in self.pages if p.text.strip())

    @property
    def word_count(self) -> int:
        return len(self.full_text.split())


class UnsupportedFileType(Exception):
    def __init__(self, ext: str) -> None:
        super().__init__(f"Unsupported file type: {ext!r}")
        self.ext = ext


def _extract_pdf(path: str) -> ExtractedDoc:
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    pages = [ExtractedPage(text=page.get_text("text"), page=i + 1) for i, page in enumerate(doc)]
    title = (doc.metadata or {}).get("title") or None
    page_count = doc.page_count
    doc.close()
    return ExtractedDoc(pages=pages, title=title, page_count=page_count)


def _extract_docx(path: str) -> ExtractedDoc:
    import docx

    document = docx.Document(path)
    text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    return ExtractedDoc(pages=[ExtractedPage(text=text)], page_count=1)


def _extract_pptx(path: str) -> ExtractedDoc:
    from pptx import Presentation

    prs = Presentation(path)
    pages: list[ExtractedPage] = []
    for i, slide in enumerate(prs.slides):
        parts = [
            shape.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text.strip()
        ]
        pages.append(ExtractedPage(text="\n".join(parts), page=i + 1))
    return ExtractedDoc(pages=pages, page_count=len(pages))


def _extract_txt(path: str) -> ExtractedDoc:
    with open(path, "r", encoding="utf-8", errors="replace") as fh:
        text = fh.read()
    return ExtractedDoc(pages=[ExtractedPage(text=text)], page_count=1)


def _extract_image(path: str) -> ExtractedDoc:
    # OCR is enabled in Phase 10 (Tesseract). Until then, images produce no text
    # rather than failing the upload.
    return ExtractedDoc(pages=[], page_count=0)


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".txt": _extract_txt,
    ".md": _extract_txt,
    ".png": _extract_image,
    ".jpg": _extract_image,
    ".jpeg": _extract_image,
}

SUPPORTED_EXTENSIONS = tuple(_EXTRACTORS.keys())


def extract(path: str, filename: str) -> ExtractedDoc:
    ext = os.path.splitext(filename)[1].lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise UnsupportedFileType(ext)
    result = extractor(path)
    if not result.title:
        result.title = os.path.splitext(os.path.basename(filename))[0]
    return result
